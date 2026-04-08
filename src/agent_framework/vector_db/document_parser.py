"""
文档解析器
支持 PDF、Word、Excel、PPT、TXT、Markdown、HTML 等多种格式
"""

import os
import re
from typing import List, Dict, Optional, Any
from dataclasses import dataclass
from abc import ABC, abstractmethod


@dataclass
class Document:
    """文档对象"""
    content: str
    metadata: Dict[str, Any]
    doc_type: str  # pdf, docx, xlsx, pptx, txt, md, html


@dataclass
class DocumentChunk:
    """文档分块"""
    content: str
    metadata: Dict[str, Any]
    chunk_index: int
    doc_id: str


class DocumentParser(ABC):
    """文档解析器基类"""

    @abstractmethod
    def parse(self, file_path: str) -> Document:
        """解析文档"""
        pass

    @abstractmethod
    def can_parse(self, file_path: str) -> bool:
        """判断是否支持该文件类型"""
        pass


class TextParser(DocumentParser):
    """纯文本解析器"""

    def can_parse(self, file_path: str) -> bool:
        return file_path.lower().endswith(('.txt', '.md', '.markdown'))

    def parse(self, file_path: str) -> Document:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()

        return Document(
            content=content,
            metadata={
                "file_name": os.path.basename(file_path),
                "file_size": os.path.getsize(file_path),
                "file_path": file_path,
            },
            doc_type="txt",
        )


class PDFParser(DocumentParser):
    """PDF 解析器"""

    def can_parse(self, file_path: str) -> bool:
        return file_path.lower().endswith('.pdf')

    def parse(self, file_path: str) -> Document:
        try:
            import PyPDF2
        except ImportError:
            raise ImportError("请安装 PyPDF2: pip install PyPDF2")

        content_parts = []
        metadata = {
            "file_name": os.path.basename(file_path),
            "file_size": os.path.getsize(file_path),
            "file_path": file_path,
        }

        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            metadata["page_count"] = len(reader.pages)

            for page_num, page in enumerate(reader.pages):
                text = page.extract_text()
                if text.strip():
                    content_parts.append(f"[Page {page_num + 1}]\n{text}")

        return Document(
            content="\n\n".join(content_parts),
            metadata=metadata,
            doc_type="pdf",
        )


class DocxParser(DocumentParser):
    """Word 文档解析器"""

    def can_parse(self, file_path: str) -> bool:
        return file_path.lower().endswith(('.docx', '.doc'))

    def parse(self, file_path: str) -> Document:
        try:
            from docx import Document as DocxDocument
        except ImportError:
            raise ImportError("请安装 python-docx: pip install python-docx")

        doc = DocxDocument(file_path)
        content_parts = []

        for para in doc.paragraphs:
            if para.text.strip():
                content_parts.append(para.text)

        # 提取表格
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells]
                table_text.append(" | ".join(row_text))
            if table_text:
                content_parts.append("\n".join(table_text))

        return Document(
            content="\n\n".join(content_parts),
            metadata={
                "file_name": os.path.basename(file_path),
                "file_size": os.path.getsize(file_path),
                "file_path": file_path,
                "paragraph_count": len(doc.paragraphs),
                "table_count": len(doc.tables),
            },
            doc_type="docx",
        )


class ExcelParser(DocumentParser):
    """Excel 解析器"""

    def can_parse(self, file_path: str) -> bool:
        return file_path.lower().endswith(('.xlsx', '.xls'))

    def parse(self, file_path: str) -> Document:
        try:
            import openpyxl
        except ImportError:
            raise ImportError("请安装 openpyxl: pip install openpyxl")

        workbook = openpyxl.load_workbook(file_path, data_only=True)
        content_parts = []

        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            content_parts.append(f"[Sheet: {sheet_name}]")

            for row in sheet.iter_rows(values_only=True):
                row_text = [str(cell) if cell is not None else "" for cell in row]
                if any(row_text):
                    content_parts.append(" | ".join(row_text))

        return Document(
            content="\n".join(content_parts),
            metadata={
                "file_name": os.path.basename(file_path),
                "file_size": os.path.getsize(file_path),
                "file_path": file_path,
                "sheet_count": len(workbook.sheetnames),
                "sheet_names": workbook.sheetnames,
            },
            doc_type="xlsx",
        )


class PPTXParser(DocumentParser):
    """PowerPoint 解析器"""

    def can_parse(self, file_path: str) -> bool:
        return file_path.lower().endswith(('.pptx', '.ppt'))

    def parse(self, file_path: str) -> Document:
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("请安装 python-pptx: pip install python-pptx")

        prs = Presentation(file_path)
        content_parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            content_parts.append(f"[Slide {slide_num}]")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    content_parts.append(shape.text)

        return Document(
            content="\n\n".join(content_parts),
            metadata={
                "file_name": os.path.basename(file_path),
                "file_size": os.path.getsize(file_path),
                "file_path": file_path,
                "slide_count": len(prs.slides),
            },
            doc_type="pptx",
        )


class HTMLParser(DocumentParser):
    """HTML 解析器"""

    def can_parse(self, file_path: str) -> bool:
        return file_path.lower().endswith(('.html', '.htm'))

    def parse(self, file_path: str) -> Document:
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            raise ImportError("请安装 beautifulsoup4: pip install beautifulsoup4")

        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            html_content = f.read()

        soup = BeautifulSoup(html_content, 'html.parser')

        # 移除 script 和 style 标签
        for script in soup(["script", "style"]):
            script.decompose()

        # 提取文本
        text = soup.get_text()

        # 清理空白
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = '\n'.join(chunk for chunk in chunks if chunk)

        return Document(
            content=text,
            metadata={
                "file_name": os.path.basename(file_path),
                "file_size": os.path.getsize(file_path),
                "file_path": file_path,
                "title": soup.title.string if soup.title else None,
            },
            doc_type="html",
        )


class DocumentLoader:
    """文档加载器 - 自动选择合适的解析器"""

    def __init__(self):
        self.parsers: List[DocumentParser] = [
            TextParser(),
            PDFParser(),
            DocxParser(),
            ExcelParser(),
            PPTXParser(),
            HTMLParser(),
        ]

    def load(self, file_path: str) -> Document:
        """加载文档"""
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在: {file_path}")

        for parser in self.parsers:
            if parser.can_parse(file_path):
                return parser.parse(file_path)

        raise ValueError(f"不支持的文件类型: {file_path}")

    def load_batch(self, file_paths: List[str]) -> List[Document]:
        """批量加载文档"""
        documents = []
        for file_path in file_paths:
            try:
                doc = self.load(file_path)
                documents.append(doc)
            except Exception as e:
                print(f"加载文件失败 {file_path}: {e}")
        return documents


class DocumentChunker:
    """文档分块器"""

    def __init__(
        self,
        chunk_size: int = 500,
        chunk_overlap: int = 50,
        separator: str = "\n\n",
    ):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.separator = separator

    def chunk(self, document: Document, doc_id: str = None) -> List[DocumentChunk]:
        """将文档分块"""
        if doc_id is None:
            doc_id = document.metadata.get("file_name", "unknown")

        # 按分隔符分割
        parts = document.content.split(self.separator)

        chunks = []
        current_chunk = []
        current_size = 0
        chunk_index = 0

        for part in parts:
            part_size = len(part)

            if current_size + part_size > self.chunk_size and current_chunk:
                # 保存当前块
                chunk_content = self.separator.join(current_chunk)
                chunks.append(DocumentChunk(
                    content=chunk_content,
                    metadata={**document.metadata, "chunk_size": len(chunk_content)},
                    chunk_index=chunk_index,
                    doc_id=doc_id,
                ))
                chunk_index += 1

                # 保留重叠部分
                if self.chunk_overlap > 0 and current_chunk:
                    overlap_text = current_chunk[-1]
                    current_chunk = [overlap_text]
                    current_size = len(overlap_text)
                else:
                    current_chunk = []
                    current_size = 0

            current_chunk.append(part)
            current_size += part_size

        # 保存最后一块
        if current_chunk:
            chunk_content = self.separator.join(current_chunk)
            chunks.append(DocumentChunk(
                content=chunk_content,
                metadata={**document.metadata, "chunk_size": len(chunk_content)},
                chunk_index=chunk_index,
                doc_id=doc_id,
            ))

        return chunks

    def chunk_by_sentences(self, document: Document, doc_id: str = None) -> List[DocumentChunk]:
        """按句子分块"""
        if doc_id is None:
            doc_id = document.metadata.get("file_name", "unknown")

        # 简单的句子分割
        sentences = re.split(r'[。！？\n]+', document.content)
        sentences = [s.strip() for s in sentences if s.strip()]

        chunks = []
        current_chunk = []
        current_size = 0
        chunk_index = 0

        for sentence in sentences:
            sentence_size = len(sentence)

            if current_size + sentence_size > self.chunk_size and current_chunk:
                chunk_content = "。".join(current_chunk) + "。"
                chunks.append(DocumentChunk(
                    content=chunk_content,
                    metadata={**document.metadata, "chunk_size": len(chunk_content)},
                    chunk_index=chunk_index,
                    doc_id=doc_id,
                ))
                chunk_index += 1
                current_chunk = []
                current_size = 0

            current_chunk.append(sentence)
            current_size += sentence_size

        if current_chunk:
            chunk_content = "。".join(current_chunk) + "。"
            chunks.append(DocumentChunk(
                content=chunk_content,
                metadata={**document.metadata, "chunk_size": len(chunk_content)},
                chunk_index=chunk_index,
                doc_id=doc_id,
            ))

        return chunks
